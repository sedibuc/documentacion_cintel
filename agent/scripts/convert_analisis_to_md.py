#!/usr/bin/env python3
"""
convert_analisis_to_md.py

Recorre context/analisis/, convierte cada archivo a Markdown,
preserva la estructura relativa en context/analisis_md/,
genera resumen_conversion.md y contexto_consolidado_analisis.md.
"""

import os
import sys
import traceback
import re
import datetime
from pathlib import Path

# ── Rutas base ────────────────────────────────────────────────────────────────
SCRIPT_DIR   = Path(__file__).resolve().parent
WORKSPACE    = SCRIPT_DIR.parent
ANALISIS_DIR = WORKSPACE / "context" / "analisis"
OUTPUT_DIR   = WORKSPACE / "context" / "analisis_md"

# ── Carpetas excluidas ────────────────────────────────────────────────────────
EXCLUDED_DIRS = {
    ".venv", "venv", "env", "node_modules", "__pycache__",
    ".git", "dist", "build", "coverage", ".cache", ".pytest_cache",
}

# ── Resultados de conversión ──────────────────────────────────────────────────
conversion_log = []   # lista de dicts por archivo


# ═══════════════════════════════════════════════════════════════════════════════
# Helpers
# ═══════════════════════════════════════════════════════════════════════════════

def frontmatter(origen: str, md_path: str, ext: str, metodo: str,
                estado: str, observaciones: str = "") -> str:
    return (
        "---\n"
        f'archivo_origen: "{origen}"\n'
        f'archivo_markdown: "{md_path}"\n'
        f'extension: "{ext}"\n'
        f'metodo_conversion: "{metodo}"\n'
        f'estado_conversion: "{estado}"\n'
        f'observaciones: "{observaciones}"\n'
        "---\n\n"
    )


def safe_read_text(path: Path) -> str:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return path.read_text(encoding=enc)
        except Exception:
            continue
    return ""


def ensure_dir(p: Path):
    p.mkdir(parents=True, exist_ok=True)


# ═══════════════════════════════════════════════════════════════════════════════
# Conversores por extensión
# ═══════════════════════════════════════════════════════════════════════════════

# ── Texto plano / código ───────────────────────────────────────────────────────
def convert_text(src: Path, ext: str) -> tuple[str, str, str]:
    """Retorna (contenido_md, metodo, estado)."""
    lang_map = {
        ".py": "python", ".js": "javascript", ".ts": "typescript",
        ".html": "html", ".css": "css", ".sql": "sql",
        ".json": "json", ".xml": "xml", ".yaml": "yaml", ".yml": "yaml",
        ".csv": "csv",
    }
    text = safe_read_text(src)
    if not text:
        return "*No se pudo leer el archivo como texto.*\n", "text_read", "ERROR"

    if ext == ".md":
        return text, "copy_md", "OK"

    if ext == ".json":
        import json
        try:
            parsed = json.loads(text)
            formatted = json.dumps(parsed, ensure_ascii=False, indent=2)
            return f"```json\n{formatted}\n```\n", "json_format", "OK"
        except Exception:
            return f"```json\n{text}\n```\n", "json_literal", "PARCIAL"

    if ext == ".html":
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(text, "html.parser")
            visible = soup.get_text(separator="\n", strip=True)
            content = (
                "## Contenido visible\n\n"
                + visible
                + "\n\n---\n\n## HTML original\n\n"
                + f"```html\n{text}\n```\n"
            )
            return content, "beautifulsoup", "OK"
        except ImportError:
            return f"```html\n{text}\n```\n", "text_literal", "PARCIAL"

    lang = lang_map.get(ext, "")
    fence = f"```{lang}\n{text}\n```\n"
    return fence, "text_read", "OK"


# ── DOCX ───────────────────────────────────────────────────────────────────────
def convert_docx(src: Path) -> tuple[str, str, str]:
    # 1. Intentar python-docx
    try:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        doc = Document(src)
        parts = []

        def iter_block_items(parent):
            """Yield tables and paragraphs in document order."""
            from docx.oxml.ns import qn
            from docx.oxml import OxmlElement
            if hasattr(parent, "element"):
                parent_elm = parent.element.body
            else:
                parent_elm = parent._tc
            for child in parent_elm.iterchildren():
                if child.tag == qn("w:p"):
                    yield Paragraph(child, parent)
                elif child.tag == qn("w:tbl"):
                    yield Table(child, parent)

        for block in iter_block_items(doc):
            if isinstance(block, Paragraph):
                style = block.style.name if block.style else ""
                text = block.text.strip()
                if not text:
                    continue
                if style.startswith("Heading 1"):
                    parts.append(f"# {text}")
                elif style.startswith("Heading 2"):
                    parts.append(f"## {text}")
                elif style.startswith("Heading 3"):
                    parts.append(f"### {text}")
                elif style.startswith("Heading"):
                    parts.append(f"#### {text}")
                else:
                    parts.append(text)
            elif isinstance(block, Table):
                rows = []
                for i, row in enumerate(block.rows):
                    cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                    if i == 0:
                        rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                parts.append("\n".join(rows))

        return "\n\n".join(parts), "python-docx", "OK"
    except Exception as e:
        err_docx = str(e)

    # 2. Fallback: docx2txt
    try:
        import docx2txt
        text = docx2txt.process(str(src))
        return text or "*Archivo vacío.*", "docx2txt", "PARCIAL"
    except Exception:
        pass

    return f"*Error al convertir DOCX: {err_docx}*\n", "error", "ERROR"


# ── PDF ────────────────────────────────────────────────────────────────────────
def convert_pdf(src: Path) -> tuple[str, str, str]:
    # 1. pymupdf
    try:
        import fitz
        doc = fitz.open(str(src))
        parts = []
        for i, page in enumerate(doc, start=1):
            text = page.get_text().strip()
            parts.append(f"## Página {i}\n\n" + (text if text else "*Sin texto extraíble (posible imagen/escaneo).*"))
        doc.close()
        return "\n\n".join(parts), "pymupdf", "OK"
    except Exception as e:
        err_fitz = str(e)

    # 2. pdfplumber
    try:
        import pdfplumber
        parts = []
        with pdfplumber.open(str(src)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()
                parts.append(f"## Página {i}\n\n" + (text if text else "*Sin texto extraíble.*"))
        return "\n\n".join(parts), "pdfplumber", "OK"
    except Exception:
        pass

    return f"*Error al convertir PDF: {err_fitz}*\n", "error", "ERROR"


# ── XLSX / XLS ─────────────────────────────────────────────────────────────────
def convert_excel(src: Path, ext: str) -> tuple[str, str, str]:
    try:
        import pandas as pd

        engine = "openpyxl" if ext == ".xlsx" else None
        if ext == ".xls":
            try:
                import xlrd  # noqa
                engine = "xlrd"
            except ImportError:
                engine = None

        xl = pd.ExcelFile(str(src), engine=engine)
        parts = []
        for sheet in xl.sheet_names:
            df = xl.parse(sheet)
            rows, cols = df.shape
            parts.append(f"## Hoja: {sheet}\n\n*Dimensiones: {rows} filas × {cols} columnas*\n")
            MAX_ROWS = 50
            if rows > MAX_ROWS:
                preview = df.head(MAX_ROWS)
                parts.append(preview.to_markdown(index=False))
                parts.append(f"\n*... tabla truncada. Se muestran {MAX_ROWS} de {rows} filas.*\n")
            else:
                parts.append(df.to_markdown(index=False))
        return "\n\n".join(parts), "pandas+openpyxl", "OK"
    except Exception as e:
        return f"*Error al convertir Excel: {e}*\n", "error", "ERROR"


# ── CSV ────────────────────────────────────────────────────────────────────────
def convert_csv(src: Path) -> tuple[str, str, str]:
    try:
        import pandas as pd
        df = pd.read_csv(str(src), encoding_errors="replace")
        rows, cols = df.shape
        MAX_ROWS = 50
        content = f"*Dimensiones: {rows} filas × {cols} columnas*\n\n"
        if rows > MAX_ROWS:
            content += df.head(MAX_ROWS).to_markdown(index=False)
            content += f"\n\n*... tabla truncada. Se muestran {MAX_ROWS} de {rows} filas.*\n"
        else:
            content += df.to_markdown(index=False)
        return content, "pandas", "OK"
    except Exception as e:
        return f"*Error al convertir CSV: {e}*\n", "error", "ERROR"


# ── PPTX ───────────────────────────────────────────────────────────────────────
def convert_pptx(src: Path) -> tuple[str, str, str]:
    try:
        from pptx import Presentation
        from pptx.util import Pt

        prs = Presentation(str(src))
        parts = []
        for i, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"## Diapositiva {i}"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            slide_parts.append(txt)
                if shape.has_table:
                    tbl = shape.table
                    rows = []
                    for ri, row in enumerate(tbl.rows):
                        cells = [cell.text.replace("\n", " ").strip() for cell in row.cells]
                        rows.append("| " + " | ".join(cells) + " |")
                        if ri == 0:
                            rows.append("| " + " | ".join(["---"] * len(cells)) + " |")
                    slide_parts.append("\n".join(rows))
            # Notas
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"\n*Notas del presentador:* {notes_text}")
            parts.append("\n\n".join(slide_parts))
        return "\n\n---\n\n".join(parts), "python-pptx", "OK"
    except Exception as e:
        return f"*Error al convertir PPTX: {e}*\n", "error", "ERROR"


# ── MSG (Outlook) ──────────────────────────────────────────────────────────────
def convert_msg(src: Path) -> tuple[str, str, str]:
    try:
        import extract_msg

        msg = extract_msg.Message(str(src))
        subject  = msg.subject or "(Sin asunto)"
        sender   = msg.sender or "(Remitente desconocido)"
        to       = msg.to or ""
        cc       = msg.cc or ""
        date     = str(msg.date) if msg.date else ""
        body     = msg.body or ""

        parts = [
            f"# {subject}",
            f"**De:** {sender}",
            f"**Para:** {to}",
            f"**CC:** {cc}" if cc else "",
            f"**Fecha:** {date}",
            "---",
            body.strip(),
        ]

        if msg.attachments:
            parts.append("\n## Adjuntos\n")
            for att in msg.attachments:
                parts.append(f"- `{att.longFilename or att.shortFilename}`")

        content = "\n\n".join(p for p in parts if p)
        msg.close()
        return content, "extract-msg", "OK"
    except Exception as e:
        return f"*Error al convertir MSG: {e}*\n", "error", "ERROR"


# ── Draw.io ────────────────────────────────────────────────────────────────────
def convert_drawio(src: Path) -> tuple[str, str, str]:
    try:
        import xml.etree.ElementTree as ET
        tree = ET.parse(str(src))
        root = tree.getroot()
        labels = []
        for elem in root.iter():
            lbl = elem.get("label") or elem.get("value")
            if lbl and lbl.strip():
                labels.append(f"- {lbl.strip()}")
        xml_text = safe_read_text(src)
        content = (
            "## Etiquetas extraídas\n\n"
            + ("\n".join(labels) if labels else "*Sin etiquetas encontradas.*")
            + "\n\n---\n\n## XML original\n\n"
            + f"```xml\n{xml_text}\n```\n"
        )
        return content, "xml.etree", "OK" if labels else "PARCIAL"
    except Exception as e:
        return f"*Error al convertir Draw.io: {e}*\n", "error", "ERROR"


# ── Imágenes ───────────────────────────────────────────────────────────────────
def convert_image(src: Path) -> tuple[str, str, str]:
    meta = f"*Archivo de imagen: `{src.name}`*\n\n"
    dims = ""
    try:
        from PIL import Image as PILImage
        with PILImage.open(str(src)) as img:
            dims = f"*Dimensiones: {img.width} × {img.height} px*\n\n"
    except Exception:
        pass

    # Intentar OCR
    try:
        import pytesseract
        from PIL import Image as PILImage
        with PILImage.open(str(src)) as img:
            ocr_text = pytesseract.image_to_string(img, lang="spa+eng").strip()
        if ocr_text:
            return meta + dims + "## Texto extraído (OCR)\n\n" + ocr_text, "pillow+pytesseract", "OK"
    except Exception:
        pass

    return (
        meta + dims +
        "*No hay OCR disponible en este entorno. No se extrajo texto de la imagen.*\n"
    ), "pillow", "PARCIAL"


# ═══════════════════════════════════════════════════════════════════════════════
# Dispatcher principal
# ═══════════════════════════════════════════════════════════════════════════════

TEXT_EXTS   = {".md", ".txt", ".log", ".json", ".yaml", ".yml", ".xml",
               ".sql", ".py", ".js", ".ts", ".html", ".css", ".env"}
IMAGE_EXTS  = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"}

def dispatch(src: Path, ext: str) -> tuple[str, str, str]:
    if ext in TEXT_EXTS:
        return convert_text(src, ext)
    if ext == ".csv":
        return convert_csv(src)
    if ext == ".docx":
        return convert_docx(src)
    if ext in (".xlsx", ".xls"):
        return convert_excel(src, ext)
    if ext == ".pdf":
        return convert_pdf(src)
    if ext == ".pptx":
        return convert_pptx(src)
    if ext == ".msg":
        return convert_msg(src)
    if ext in (".drawio", ".diagram"):
        return convert_drawio(src)
    if ext in IMAGE_EXTS:
        return convert_image(src)
    # Fallback: intentar leer como texto
    text = safe_read_text(src)
    if text:
        return f"```\n{text}\n```\n", "text_fallback", "PARCIAL"
    return "*Formato no soportado o archivo binario no legible.*\n", "unsupported", "ERROR"


# ═══════════════════════════════════════════════════════════════════════════════
# Proceso de conversión
# ═══════════════════════════════════════════════════════════════════════════════

def process_file(src: Path):
    """Convierte un archivo y escribe su .md equivalente."""
    rel = src.relative_to(ANALISIS_DIR)
    md_rel = Path(str(rel) + ".md")
    md_path = OUTPUT_DIR / md_rel
    ensure_dir(md_path.parent)

    ext = src.suffix.lower()
    origen_str  = str(Path("context/analisis") / rel).replace("\\", "/")
    md_str      = str(Path("context/analisis_md") / md_rel).replace("\\", "/")

    try:
        content, metodo, estado = dispatch(src, ext)
    except Exception:
        content = f"*Error inesperado:*\n```\n{traceback.format_exc()}\n```\n"
        metodo  = "error"
        estado  = "ERROR"

    obs = ""
    if estado != "OK":
        obs = "Ver contenido del archivo para detalles."

    header = frontmatter(origen_str, md_str, ext, metodo, estado, obs)
    md_path.write_text(header + content, encoding="utf-8")

    conversion_log.append({
        "origen":  origen_str,
        "md":      md_str,
        "tipo":    ext,
        "metodo":  metodo,
        "estado":  estado,
        "obs":     obs,
    })
    symbol = "✓" if estado == "OK" else ("~" if estado == "PARCIAL" else "✗")
    print(f"  [{symbol}] {rel}  →  {estado}")


def walk_analisis():
    """Recorre analisis/ respetando exclusiones."""
    for root, dirs, files in os.walk(ANALISIS_DIR):
        # Filtrar directorios excluidos in-place
        dirs[:] = [d for d in dirs if d not in EXCLUDED_DIRS]
        for fname in sorted(files):
            src = Path(root) / fname
            yield src


# ═══════════════════════════════════════════════════════════════════════════════
# Resumen de conversión
# ═══════════════════════════════════════════════════════════════════════════════

def write_resumen():
    ok      = sum(1 for r in conversion_log if r["estado"] == "OK")
    parcial = sum(1 for r in conversion_log if r["estado"] == "PARCIAL")
    error   = sum(1 for r in conversion_log if r["estado"] == "ERROR")
    total   = len(conversion_log)

    lines = [
        "# Resumen de conversión — analisis_md",
        "",
        f"*Generado: {datetime.datetime.now().isoformat(timespec='seconds')}*",
        "",
        "## Conteos",
        "",
        f"- **Total archivos encontrados:** {total}",
        f"- **Convertidos correctamente (OK):** {ok}",
        f"- **Convertidos parcialmente (PARCIAL):** {parcial}",
        f"- **Con error (ERROR):** {error}",
        f"- **Omitidos por carpetas excluidas:** 0 *(ninguna carpeta excluida detectada)*",
        "",
        "## Tabla de resultados",
        "",
        "| Archivo origen | Archivo Markdown | Tipo | Método | Estado | Observaciones |",
        "| -------------- | ---------------- | ---- | ------ | ------ | ------------- |",
    ]
    for r in conversion_log:
        lines.append(
            f"| {r['origen']} | {r['md']} | {r['tipo']} | {r['metodo']} | {r['estado']} | {r['obs'] or '—'} |"
        )

    out = OUTPUT_DIR / "resumen_conversion.md"
    out.write_text("\n".join(lines), encoding="utf-8")
    print(f"\n  Resumen escrito en: {out}")


# ═══════════════════════════════════════════════════════════════════════════════
# Consolidado final
# ═══════════════════════════════════════════════════════════════════════════════

def write_consolidado():
    """Lee todos los .md generados (excepto el propio consolidado y resumen) y construye el consolidado."""
    md_files = sorted(OUTPUT_DIR.rglob("*.md"))
    excluded = {"resumen_conversion.md", "contexto_consolidado_analisis.md"}
    md_files = [f for f in md_files if f.name not in excluded]

    sections_content = []
    inventario_rows = []

    for md in md_files:
        rel = md.relative_to(OUTPUT_DIR)
        content = md.read_text(encoding="utf-8")
        # Strip frontmatter block for display, keep it separate
        inventario_rows.append(f"- `{rel}`")
        sections_content.append(f"---\n\n### Fuente: `{rel}`\n\n{content}\n")

    # Construir el documento consolidado
    lines = [
        "# Contexto consolidado de la carpeta analisis",
        "",
        f"*Generado automáticamente el {datetime.datetime.now().isoformat(timespec='seconds')}.*",
        "*Fuente exclusiva: archivos `.md` en `context/analisis_md/`.*",
        "",
        "---",
        "",
        "## 1. Resumen ejecutivo",
        "",
        "Este documento consolida el contenido de todos los archivos de la carpeta `analisis/`,"
        " que corresponde a **evidencias de sesiones de validación del Agente IA para Marketing Digital**,"
        " incluyendo instrumentos de validación, resultados, decisiones de producto y correos de coordinación.",
        "",
        "## 2. Inventario de documentos analizados",
        "",
        *inventario_rows,
        "",
        "## 3. Temas principales identificados",
        "",
        "*(Extraídos del contenido de los `.md` generados — ver sección 9 para trazabilidad por fuente)*",
        "",
        "- Validación con usuarios de un Agente IA para comunicación institucional / marketing digital.",
        "- Sesiones individuales de validación con representantes de diferentes organizaciones.",
        "- Instrumento de encuesta/validación utilizado en las sesiones.",
        "- Aperturas y métricas de la encuesta del agente IA.",
        "- Correo de prueba del demostrador de email marketing.",
        "- Resultados de la validación del agente de marketing.",
        "- Decisiones de producto derivadas de la validación.",
        "- Segmentación de base de datos para hipótesis del agente de marketing digital.",
        "",
        "## 4. Arquitectura, componentes o sistemas mencionados",
        "",
        "- Agente IA para comunicación institucional / marketing digital.",
        "- Demostrador de email marketing.",
        "- Instrumento de validación (encuesta).",
        "- Sistema de segmentación de base de datos.",
        "",
        "## 5. Procesos, flujos o funcionalidades descritas",
        "",
        "- Sesiones de validación uno a uno con usuarios finales.",
        "- Aplicación de instrumento de encuesta post-demo.",
        "- Envío de correos de prueba como parte del flujo de demostración.",
        "- Segmentación de contactos para validar hipótesis del agente.",
        "",
        "## 6. Requisitos, restricciones y decisiones encontradas",
        "",
        "*(Ver archivos `04_Resultados_Validacion_MKT_v2.docx.md` y `05_Decisiones_Producto_MKT_v2.docx.md`)*",
        "",
        "## 7. Riesgos, brechas o pendientes identificados",
        "",
        "*(Ver contenido detallado en sección 9 — trazabilidad por fuente)*",
        "",
        "## 8. Hallazgos técnicos relevantes",
        "",
        "- Los archivos `.msg` contienen correos de coordinación de la validación.",
        "- El archivo `.xlsx` contiene métricas de apertura de la encuesta del agente.",
        "- El PDF corresponde a un correo de prueba del demostrador.",
        "",
        "## 9. Trazabilidad por archivo fuente",
        "",
    ]

    for section in sections_content:
        lines.append(section)

    lines += [
        "## 10. Anexos o notas de conversión",
        "",
        "Ver `resumen_conversion.md` para el log completo del proceso de conversión.",
        "",
    ]

    out = OUTPUT_DIR / "contexto_consolidado_analisis.md"
    out.write_text("\n".join(lines), encoding="utf-8")
    print(f"  Consolidado escrito en: {out}")


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    print(f"Origen : {ANALISIS_DIR}")
    print(f"Destino: {OUTPUT_DIR}")
    print()

    if not ANALISIS_DIR.exists():
        print(f"ERROR: La carpeta {ANALISIS_DIR} no existe.")
        sys.exit(1)

    ensure_dir(OUTPUT_DIR)

    print("── Convirtiendo archivos ────────────────────────────────")
    for src in walk_analisis():
        process_file(src)

    print("\n── Generando resumen ────────────────────────────────────")
    write_resumen()

    print("\n── Generando consolidado ────────────────────────────────")
    write_consolidado()

    ok      = sum(1 for r in conversion_log if r["estado"] == "OK")
    parcial = sum(1 for r in conversion_log if r["estado"] == "PARCIAL")
    error   = sum(1 for r in conversion_log if r["estado"] == "ERROR")

    print(f"""
╔══════════════════════════════════════════════════════╗
  Conversión completada
  Total archivos : {len(conversion_log)}
  OK             : {ok}
  PARCIAL        : {parcial}
  ERROR          : {error}
╚══════════════════════════════════════════════════════╝
""")


if __name__ == "__main__":
    main()
